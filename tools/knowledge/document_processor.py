"""文档解析与分块 —— 支持 txt/md/pdf/docx/pptx/xlsx(可选 doc/ppt 转换)。"""

from __future__ import annotations

import re
import unicodedata
from dataclasses import dataclass, field
from pathlib import Path
from typing import List

from .source_role import classify_source_role, heading_level

SUPPORTED_EXTS = {".txt", ".md", ".pdf", ".docx", ".pptx", ".xlsx"}

# Excel 分块参数
EXCEL_ROWS_PER_CHUNK = 30
EXCEL_HEADER_DETECT_MAX_ROWS = 5
EXCEL_DATA_ROW_MIN_NUMERIC = 2
EXCEL_DATA_ROW_NUMERIC_RATIO = 0.5

_OCR_NOISE_RE = re.compile(
    r"^(?:NOTEBOOK|NO\.?|No\.?|DATE|Date|日期|页码|第\s*\d+\s*页|Page\s*\d+)\s*[:：]?\s*\w*\s*$",
    re.I,
)
_TITLE_KEYWORDS = ("定义", "性质", "定理", "规则", "方法", "公式", "例题", "易错", "注意", "总结", "步骤")
_BODY_LIKE_ENDINGS = ("。", "；", ";", ".", "！", "？", "!", "?")
_FORMULA_RE = re.compile(
    r"(?:\\\[|\\\(|\$\$?|[A-Za-z][A-Za-z0-9_']*\s*=|[∑√≈≠≤≥→←⇔]|[+\-*/=]\s*[A-Za-z0-9(])"
)
_DISPLAY_FORMULA_RE = re.compile(r"(\$\$.*?\$\$|\\\[.*?\\\]|\\\(.*?\\\))", re.S)
_VISUAL_REGION_RE = re.compile(r"<!--\s*图示:\s*([^;>-]+).*?type=([a-zA-Z0-9_]+)")
_CONTENT_TAG_RULES = (
    ("definition", re.compile(r"(定义|概念|称为)")),
    ("theorem", re.compile(r"(定理|性质|规律|规则)")),
    ("formula", re.compile(r"(公式|计算式|表达式|方程|等式)")),
    ("method", re.compile(r"(方法|步骤|解法|技巧|化简|证明)")),
    ("example", re.compile(r"(例题|例\s*\d*|example)", re.I)),
    ("mistake", re.compile(r"(易错|注意|误区|不要|错因|陷阱)")),
)


@dataclass
class TextChunk:
    """单个文本块"""
    text: str
    metadata: dict = field(default_factory=dict)


# ============================================================
# 文本清洗(减少 embedding NaN / 噪声)
# ============================================================
def sanitize_text(text: str, keep_newlines: bool = False) -> str:
    """NFKC 归一化 + 剔除控制字符 + 压缩空白。keep_newlines 时按行清洗，留给标题切块。"""
    if not isinstance(text, str):
        return str(text) if text is not None else ""
    text = unicodedata.normalize("NFKC", text)
    text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]", "", text)
    if keep_newlines:
        lines = []
        for line in text.replace("\r\n", "\n").replace("\r", "\n").split("\n"):
            line = re.sub(r"[ \t]+", " ", line).strip()
            if line:
                lines.append(line)
        return "\n".join(lines)
    text = re.sub(r"\s+", " ", text)
    return text.strip()


def _base_meta(path: str, text: str = "", role: str = "", **extra: object) -> dict:
    name = Path(path).name
    meta = {"source": name, "role": role or classify_source_role(name, text)}
    meta.update({key: value for key, value in extra.items() if value not in (None, "")})
    return meta


def _clean_ocr_markdown(text: str) -> str:
    """清理 OCR Markdown 常见页眉页脚、重复空行和相邻重复标题。"""
    lines = sanitize_text(text, keep_newlines=True).splitlines()
    out: list[str] = []
    last_heading = ""
    for line in lines:
        if _OCR_NOISE_RE.match(line):
            continue
        if re.fullmatch(r"[-_=]{3,}", line):
            continue
        hit = heading_level(line)
        if hit:
            level, title = hit
            key = f"{level}:{title}"
            if key == last_heading:
                continue
            last_heading = key
        elif line:
            last_heading = ""
        out.append(line)
    return "\n".join(out)


def _heading_score(title: str, level: int, path_len: int) -> int:
    title = " ".join(str(title or "").split()).strip()
    score = 0
    if level == 1:
        score += 5
    elif level == 2:
        score += 4
    else:
        score += 2
    if path_len >= 2:
        score += 2
    if any(word in title for word in _TITLE_KEYWORDS):
        score += 2
    if 2 <= len(title) <= 28:
        score += 1
    if len(title) > 45:
        score -= 3
    if title.endswith(_BODY_LIKE_ENDINGS):
        score -= 3
    return max(score, 0)


def _heading_kind(title: str, level: int, score: int) -> str:
    title = str(title or "")
    if score < 4:
        return "evidence"
    if level == 1:
        return "chapter"
    if level == 2:
        return "topic"
    if any(word in title for word in ("易错", "注意", "例题", "例")):
        return "evidence"
    return "knowledge_point"


def _content_tags(text: str, heading: str = "") -> str:
    blob = f"{heading}\n{text or ''}"
    tags = [name for name, pattern in _CONTENT_TAG_RULES if pattern.search(blob)]
    if _FORMULA_RE.search(blob) and "formula" not in tags:
        tags.append("formula")
    if _VISUAL_REGION_RE.search(blob):
        tags.append("visual")
    return ",".join(tags)


def _visual_region_meta(text: str) -> dict[str, object]:
    hits = _VISUAL_REGION_RE.findall(text or "")
    if not hits:
        return {}
    types = sorted({str(kind).strip() for _label, kind in hits if str(kind).strip()})
    labels = sorted({str(label).strip() for label, _kind in hits if str(label).strip()})
    return {
        "visual_region_count": len(hits),
        "visual_region_types": ",".join(types),
        "visual_region_labels": "、".join(labels),
    }


def _contains_formula(text: str) -> bool:
    return bool(_FORMULA_RE.search(text or ""))


def _chunks_by_heading(text: str, path: str) -> List[TextChunk]:
    """按标题切开，块上带完整 heading_path / score / kind / content_tags。"""
    name = Path(path).name
    text = _clean_ocr_markdown(text)
    role = classify_source_role(name, text)
    heading_stack: dict[int, str] = {}
    current_level = 0
    buf: list[str] = []
    chunks: List[TextChunk] = []

    def flush() -> None:
        body = "\n".join(buf).strip()
        buf.clear()
        if not body:
            return
        path_parts = [heading_stack[i] for i in sorted(heading_stack) if heading_stack.get(i)]
        heading = path_parts[-1] if path_parts else name
        level = current_level or len(path_parts) or 1
        score = _heading_score(heading, level, len(path_parts))
        kind = _heading_kind(heading, level, score)
        chunks.append(
            TextChunk(
                sanitize_text(body, keep_newlines=True),
                _base_meta(
                    path,
                    role=role,
                    chapter=path_parts[0] if path_parts else "",
                    topic=path_parts[1] if len(path_parts) > 1 else "",
                    heading=heading,
                    heading_level=level,
                    heading_depth=level,
                    heading_path_text=" / ".join(path_parts),
                    heading_score=score,
                    heading_kind=kind,
                    content_tags=_content_tags(body, heading),
                    contains_formula=_contains_formula(body),
                    block_type="formula_heavy" if _contains_formula(body) else "content",
                    **_visual_region_meta(body),
                ),
            )
        )

    for line in text.splitlines():
        hit = heading_level(line)
        if hit:
            flush()
            level, title = hit
            current_level = level
            heading_stack[level] = title
            for old_level in list(heading_stack):
                if old_level > level:
                    heading_stack.pop(old_level, None)
            continue
        buf.append(line)
    flush()
    if chunks:
        return chunks
    body = sanitize_text(text, keep_newlines=True)
    if not body:
        return []
    return [
        TextChunk(
            body,
            _base_meta(
                path,
                role=role,
                heading=name,
                heading_level=0,
                heading_depth=0,
                heading_path_text=name,
                heading_score=0,
                heading_kind="evidence",
                content_tags=_content_tags(body, name),
                contains_formula=_contains_formula(body),
                block_type="formula_heavy" if _contains_formula(body) else "content",
                **_visual_region_meta(body),
            ),
        )
    ]


# ============================================================
# 各格式提取
# ============================================================
def _extract_txt(path: str) -> List[TextChunk]:
    for enc in ("utf-8", "gbk", "latin-1"):
        try:
            with open(path, "r", encoding=enc) as f:
                raw = f.read()
            break
        except (UnicodeDecodeError, LookupError):
            continue
    else:
        raw = Path(path).read_text(errors="ignore")
    return _chunks_by_heading(raw, path)


def _extract_pdf(path: str) -> List[TextChunk]:
    import pdfplumber
    chunks = []
    with pdfplumber.open(path) as pdf:
        pages_text = [page.extract_text() or "" for page in pdf.pages]
        role = classify_source_role(Path(path).name, "\n".join(pages_text))
        for i, text in enumerate(pages_text):
            if text.strip():
                first = sanitize_text(text, keep_newlines=True).splitlines()
                heading = first[0] if first else ""
                chunks.append(TextChunk(
                    sanitize_text(text, keep_newlines=True),
                    _base_meta(path, role=role, page=i + 1, heading=heading)))
    return chunks


def _extract_docx(path: str) -> List[TextChunk]:
    import docx
    doc = docx.Document(path)
    full_text = "\n".join(p.text.strip() for p in doc.paragraphs if p.text.strip())
    role = classify_source_role(Path(path).name, full_text)
    chapter = ""
    topic = ""
    heading = ""
    buf: list[str] = []
    chunks: List[TextChunk] = []

    def flush() -> None:
        body = "\n".join(buf).strip()
        buf.clear()
        if not body:
            return
        chunks.append(
            TextChunk(
                sanitize_text(body, keep_newlines=True),
                _base_meta(
                    path,
                    role=role,
                    chapter=chapter,
                    topic=topic,
                    heading=heading or topic or chapter,
                ),
            )
        )

    for para in doc.paragraphs:
        raw = (para.text or "").strip()
        if not raw:
            continue
        style = str(getattr(para.style, "name", "") or "").lower()
        if "heading" in style:
            flush()
            heading = raw
            if "1" in style:
                chapter = raw
                topic = ""
            elif "2" in style:
                topic = raw
            continue
        hit = heading_level(raw)
        if hit:
            flush()
            level, title = hit
            heading = title
            if level == 1:
                chapter = title
                topic = ""
            elif level == 2:
                topic = title
            continue
        buf.append(raw)
    flush()
    if chunks:
        return chunks
    return _chunks_by_heading(full_text, path)


def _extract_pptx(path: str) -> List[TextChunk]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    chunks = []
    prs = Presentation(path)
    slides_parts: list[list[str]] = []
    for i, slide in enumerate(prs.slides):
        parts = []
        for shape in _iter_shapes(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # 图片内容暂不解析，直接跳过
                continue
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    cells = " | ".join(c.text.strip() for c in row.cells)
                    if cells.strip():
                        parts.append(cells)
        slides_parts.append(parts)
    role = classify_source_role(Path(path).name, "\n".join("\n".join(parts) for parts in slides_parts))
    for i, parts in enumerate(slides_parts):
        if parts:
            heading = sanitize_text(parts[0])
            chunks.append(TextChunk(
                sanitize_text("\n".join(parts), keep_newlines=True),
                _base_meta(path, role=role, page=i + 1, heading=heading)))
    return chunks


def _iter_shapes(shapes):
    """递归遍历形状(含组合形状 group)"""
    for shape in shapes:
        yield shape
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from _iter_shapes(shape.shapes)
        except Exception:
            pass


# ============================================================
# Excel: 表头识别 + 统计摘要
# ============================================================
def _is_numeric_cell(val) -> bool:
    if isinstance(val, (int, float)):
        return True
    if isinstance(val, str):
        s = val.strip()
        if re.match(r"^\d{4}[-/]\d{1,2}[-/]\d{1,2}", s):
            return True
        try:
            float(s)
            return True
        except ValueError:
            return False
    return False


def _detect_header_rows(rows_raw) -> int:
    """启发式识别表头行数: 找到第一个「像数据行」的行号(至少 1, 最多 5)"""
    for idx, row in enumerate(rows_raw[:EXCEL_HEADER_DETECT_MAX_ROWS]):
        cells = [str(c).strip() for c in row if c is not None]
        if not cells:
            continue
        numeric = sum(1 for c in cells if _is_numeric_cell(c))
        if numeric >= EXCEL_DATA_ROW_MIN_NUMERIC and \
           numeric / len(cells) >= EXCEL_DATA_ROW_NUMERIC_RATIO:
            return max(idx, 1)
    return 1


def _excel_summary(base, sheet, sheet_idx, rows_raw, header_rows) -> TextChunk:
    """统计摘要块: 总行数/列名/每列统计/数值列 合计平均最大最小"""
    lines = [f"【统计摘要】文件: {base}，工作表: {sheet}",
             f"总数据行数: {max(len(rows_raw) - header_rows, 0)} 行（不含表头）"]
    data_rows = rows_raw[header_rows:]
    if rows_raw:
        ncol = max((len(r) for r in rows_raw), default=0)
        lines.append(f"列数: {ncol}")
        header = [str(c).strip() for c in rows_raw[0][:ncol] if c is not None]
        if header:
            lines.append("表头列名: " + ", ".join(header))
        for col in range(ncol):
            vals = [r[col] for r in data_rows if col < len(r)]
            non_empty = [v for v in vals if v is not None and str(v).strip() != ""]
            lines.append(f"列{col+1}: {len(non_empty)}个非空值, "
                         f"{len(set(str(v) for v in non_empty))}个不同值, "
                         f"{len(vals) - len(non_empty)}个空值")
            if non_empty and all(_is_numeric_cell(v) for v in non_empty):
                nums = [float(v) for v in non_empty]
                lines.append(f"    合计={sum(nums):.2f}, 平均={sum(nums)/len(nums):.2f}, "
                             f"最小={min(nums):.2f}, 最大={max(nums):.2f}")
    return TextChunk("\n".join(lines),
                     {"source": base, "sheet": sheet, "sheet_idx": sheet_idx,
                      "type": "summary"})


def _extract_xlsx(path: str) -> List[TextChunk]:
    import openpyxl
    wb = openpyxl.load_workbook(path, data_only=True)
    chunks = []
    for si, ws in enumerate(wb.worksheets):
        rows_raw = [[c.value for c in row] for row in ws.iter_rows()]
        rows_raw = [r for r in rows_raw if any(c is not None for c in r)]
        if not rows_raw:
            continue
        header_rows = _detect_header_rows(rows_raw)
        # 表头 + 每 N 行一块
        header = " | ".join(str(c).strip() if c is not None else "" for c in rows_raw[0])
        chunks.append(TextChunk(header, {"source": Path(path).name,
                                         "sheet": ws.title, "header": True}))
        for i in range(header_rows, len(rows_raw), EXCEL_ROWS_PER_CHUNK):
            block_rows = rows_raw[i:i + EXCEL_ROWS_PER_CHUNK]
            block = "\n".join(
                " | ".join(str(c).strip() if c is not None else "" for c in r)
                for r in block_rows)
            if block.strip():
                chunks.append(TextChunk(sanitize_text(block),
                                        {"source": Path(path).name,
                                         "sheet": ws.title,
                                         "rows": f"{i-header_rows+1}-{i+len(block_rows)-header_rows}"}))
        chunks.append(_excel_summary(Path(path).name, ws.title, si, rows_raw, header_rows))
    return chunks


# ============================================================
# 分块
# ============================================================
def _estimated_tokens(text: str) -> int:
    """粗估 token 数：中文约 1 字 1 token，英文/符号约 4 字 1 token，LaTeX 公式折算更低。"""
    if not text:
        return 0
    total = 0.0
    pos = 0
    for match in _DISPLAY_FORMULA_RE.finditer(text):
        if match.start() > pos:
            total += _estimated_tokens_plain(text[pos:match.start()])
        total += max(1.0, len(match.group(0)) * 0.35)
        pos = match.end()
    if pos < len(text):
        total += _estimated_tokens_plain(text[pos:])
    return max(1, int(total + 0.999))


def _estimated_tokens_plain(text: str) -> float:
    cjk = sum(1 for ch in text if "\u4e00" <= ch <= "\u9fff")
    other = max(0, len(text) - cjk)
    return cjk + other / 4


def _split_units(text: str, separators: List[str]) -> list[tuple[str, bool]]:
    """切成可合并单元；公式片段作为不可拆单元。"""
    units: list[tuple[str, bool]] = []
    pos = 0
    for match in _DISPLAY_FORMULA_RE.finditer(text):
        if match.start() > pos:
            units.extend(_split_plain_units(text[pos:match.start()], separators))
        units.append((match.group(0), True))
        pos = match.end()
    if pos < len(text):
        units.extend(_split_plain_units(text[pos:], separators))
    return [(unit, protected) for unit, protected in units if unit]


def _split_plain_units(text: str, separators: List[str]) -> list[tuple[str, bool]]:
    units = [text]
    for sep in separators:
        next_units: list[str] = []
        for unit in units:
            if sep not in unit:
                next_units.append(unit)
                continue
            parts = unit.split(sep)
            for idx, part in enumerate(parts):
                if part:
                    next_units.append(part)
                if idx < len(parts) - 1:
                    next_units.append(sep)
        units = next_units
    return [(unit, False) for unit in units if unit]


def split_text(chunks: List[TextChunk], chunk_size: int, chunk_overlap: int) -> List[TextChunk]:
    """把超大块按估算 token 切分；公式片段尽量保持完整。"""
    result: List[TextChunk] = []
    for ci, chunk in enumerate(chunks):
        if _estimated_tokens(chunk.text) <= chunk_size:
            meta = dict(chunk.metadata)
            meta["estimated_tokens"] = _estimated_tokens(chunk.text)
            meta["split_strategy"] = "token_aware"
            if _DISPLAY_FORMULA_RE.search(chunk.text) and meta["estimated_tokens"] > chunk_size:
                meta["oversized_formula"] = True
            chunk.metadata = meta
            result.append(chunk)
            continue
        pieces = _token_aware_split(
            chunk.text,
            ["\n\n", "\n", "。", "；", "，", " "],
            chunk_size,
            chunk_overlap,
        )
        for pi, piece in enumerate(pieces):
            piece = sanitize_text(piece)
            if len(piece) < 30:
                continue
            meta = dict(chunk.metadata)
            meta["chunk_index"] = f"{ci}-{pi}"
            meta["estimated_tokens"] = _estimated_tokens(piece)
            meta["split_strategy"] = "token_aware"
            if _DISPLAY_FORMULA_RE.search(piece) and meta["estimated_tokens"] > chunk_size:
                meta["oversized_formula"] = True
            result.append(TextChunk(piece, meta))
    return result


def _token_aware_split(text: str, separators: List[str], chunk_size: int, overlap: int) -> List[str]:
    if _estimated_tokens(text) <= chunk_size:
        return [text]
    units = _split_units(text, separators)
    if not units:
        return []
    pieces: list[str] = []
    buf: list[str] = []
    buf_tokens = 0
    overlap_units: list[str] = []
    overlap_tokens = 0

    def push_buf() -> None:
        nonlocal buf, buf_tokens, overlap_units, overlap_tokens
        body = "".join(buf).strip()
        if body:
            pieces.append(body)
        if overlap > 0:
            overlap_units = []
            overlap_tokens = 0
            for unit in reversed(buf):
                tokens = _estimated_tokens(unit)
                if overlap_units and overlap_tokens + tokens > overlap:
                    break
                overlap_units.insert(0, unit)
                overlap_tokens += tokens
        buf = list(overlap_units)
        buf_tokens = overlap_tokens

    for unit, protected in units:
        tokens = _estimated_tokens(unit)
        if not protected and tokens > chunk_size:
            if buf:
                push_buf()
            pieces.extend(_hard_split_plain(unit, chunk_size, overlap))
            buf = []
            buf_tokens = 0
            overlap_units = []
            overlap_tokens = 0
            continue
        if not buf or buf_tokens + tokens <= chunk_size:
            buf.append(unit)
            buf_tokens += tokens
        else:
            push_buf()
            buf.append(unit)
            buf_tokens += tokens
    if buf:
        body = "".join(buf).strip()
        if body and (not pieces or body != pieces[-1]):
            pieces.append(body)
    return pieces


def _hard_split_plain(text: str, chunk_size: int, overlap: int) -> list[str]:
    step_tokens = max(chunk_size - overlap, 1)
    pieces: list[str] = []
    buf = ""
    buf_tokens = 0
    idx = 0
    while idx < len(text):
        ch = text[idx]
        ch_tokens = _estimated_tokens_plain(ch)
        if buf and buf_tokens + ch_tokens > chunk_size:
            pieces.append(buf.strip())
            keep = ""
            keep_tokens = 0
            for old_ch in reversed(buf):
                tok = _estimated_tokens_plain(old_ch)
                if keep and keep_tokens + tok > overlap:
                    break
                keep = old_ch + keep
                keep_tokens += tok
            buf = keep
            buf_tokens = keep_tokens
            if buf_tokens >= step_tokens:
                buf = ""
                buf_tokens = 0
            continue
        buf += ch
        buf_tokens += ch_tokens
        idx += 1
    if buf.strip():
        pieces.append(buf.strip())
    return pieces


# ============================================================
# 总入口
# ============================================================
def process_file(file_path: str, chunk_size: int = 500,
                 chunk_overlap: int = 100) -> List[TextChunk]:
    """解析单个文件并分块。Excel 已按行组织, 不再二次切分。"""
    path = Path(file_path)
    if not path.is_file():
        raise FileNotFoundError(f"文件不存在: {file_path}")
    ext = path.suffix.lower()
    if ext not in SUPPORTED_EXTS:
        raise ValueError(
            f"不支持的文件格式 '{ext}'，支持: {', '.join(sorted(SUPPORTED_EXTS))}")

    handlers = {
        ".txt": _extract_txt,
        ".md": _extract_txt,
        ".pdf": _extract_pdf,
        ".docx": _extract_docx,
        ".pptx": _extract_pptx,
        ".xlsx": _extract_xlsx,
    }
    chunks = handlers[ext](str(path))
    if ext in (".xlsx",):
        return chunks                       # Excel 不再 split_text
    return split_text(chunks, chunk_size, chunk_overlap)
